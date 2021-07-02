# -*- coding: UTF-8 -*-
# ! /usr/bin/python3

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
import copy
import six
import os
import setting as config


class PPTX(object):
    def __init__(self):
        self.template = None
        self.added_pptx = None
        self.__get_template()

    def __get_template(self):
        self.template = Presentation(config.PPTX_TEMPLATE)

    def save(self, filepath):
        save_filepath = self.__get_save_filepath(filepath)
        self.added_pptx.save(save_filepath)
        return save_filepath

    def __get_save_filepath(self, filepath):
        if os.path.exists(filepath):
            filepath_map = os.path.splitext(filepath)
            return self.__get_save_filepath(filepath_map[0] + '(0)' + filepath_map[1])
        else:
            return filepath

    def copy_slide(self, rows_count):
        f"""
         复制 PPT 幻灯片
        :param rows_count: {int}
        :return: 
        """
        try:
            self.added_pptx = self.template

            for index in range(1, rows_count):
                # 获取幻灯片
                slide = self.added_pptx.slides[0]
                blank_slide_layout = self.added_pptx.slide_layouts[0]
                # 追加幻灯片
                copy_slide = self.added_pptx.slides.add_slide(blank_slide_layout)
                shapes = slide.shapes

                # 遍历幻灯片中的形状 Shape 数组
                for shape_item in shapes:
                    el = shape_item.element
                    new_el = copy.deepcopy(el)
                    copy_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

                for _, value in six.iteritems(slide.part.rels):
                    # Make sure we don't copy a notesSlide relation as that won't exist
                    if "notesSlide" not in value.reltype:
                        copy_slide.part.rels.add_relationship(value.reltype, value._target, value.rId)
                print('添加幻灯片第 %s 页' % index)
            # 移除模板 PPT

            # for item in self.added_pptx.slides._sldIdLst:
            #     print(item.rId, type(self.added_pptx.slides._sldIdLst))
            del self.added_pptx.slides._sldIdLst[0]
            # r_id = self.added_pptx.slides._sldIdLst[-1].rId
            # self.added_pptx.part.drop_rel(r_id)
            # del self.added_pptx.slides._sldIdLst[1]
            # self.added_pptx.slides.part.drop_rel(1)

            print('幻灯片总页数为%s' % len(self.added_pptx.slides))

        except ValueError:
            pass

    def set_data_with_slide(self, **kw):
        f"""
        设置单个幻灯片
        :param kw: 
        :return: 
        """
        # :param data: {(data"case_code", "rescuer", "rescuer_phone", "rescuer_place", "use_exp_date")}
        # # 当前幻灯片
        print('更新幻灯片第 %s 页: %s ' % (kw['index'], kw))
        current_slide = self.added_pptx.slides[kw['index'] - 1]

        case_code = kw['case_code']
        rescuer = kw['rescuer']
        rescuer_phone = kw['rescuer_phone']
        rescuer_place = kw['rescuer_place']
        use_exp_date = kw['use_exp_date']

        # PPT 中生成券号文案
        s_thead = current_slide.shapes[6].text_frame
        s_thead.text = case_code if case_code else ''
        s_thead.margin_top = Pt(10)
        s_thead.paragraphs[0].font.bold = True
        s_thead.paragraphs[0].font.size = Pt(70)
        s_thead.paragraphs[0].font.color.rgb = RGBColor(127, 127, 127)

        # PPT 中生成申请人文案，格式为 申请人： 救助人真实姓名 + ' ' + 联系电话
        s_thead = current_slide.shapes[3].text_frame
        s_thead.paragraphs[0].font.name = 'Microsoft YaHei'
        # s_thead.text = '申请人：' + rescuer if rescuer is None else '' + ' ' +
        # rescuer_phone if rescuer_phone is None else ''
        s_thead.text = '申请人：%s %s' % (
            rescuer if rescuer else '', rescuer_phone if rescuer_phone else '')
        s_thead.paragraphs[0].font.color.rgb = RGBColor(127, 127, 127)
        s_thead.paragraphs[0].font.size = Pt(42)

        # PPT 中生成申请医院文案，格式为 申请医院： + 医院名称
        s_thead = current_slide.shapes[4].text_frame
        s_thead.paragraphs[0].font.name = 'Microsoft YaHei'
        s_thead.text = '申请医院：%s' % rescuer_place if rescuer_place else ''
        # s_thead.text = '申请医院：' + rescuer_place if rescuer_place else ''
        s_thead.paragraphs[0].font.color.rgb = RGBColor(127, 127, 127)
        s_thead.paragraphs[0].font.size = Pt(42)

        # PPT 中生成有效期文案，格式为 使用有效期： + 有效期
        s_thead = current_slide.shapes[5].text_frame
        s_thead.text = '使用有效期：%s' % use_exp_date if use_exp_date else ''
        # s_thead.text = '使用有效期：' + use_exp_date if use_exp_date else ''
        s_thead.paragraphs[0].font.color.rgb = RGBColor(127, 127, 127)
        s_thead.paragraphs[0].font.size = Pt(42)
